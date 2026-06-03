"""Document text extraction for SPEC_014 (user document upload + NER).

Inspired by Proto_Demo's FormattingExtractor (`Proto_Demo/src/formatter/extractor.py`)
but simplified for our NER use case — we return plain text per page rather than
the full FormattedDocument IR (paragraphs / lines / spans / fonts / colors).

Supported formats:
  - PDF (via pdfplumber — already a Market Zero dependency); per-page text + tables
  - PPTX (via python-pptx); per-slide text + tables + speaker notes — DR-9, the
    conference-deck (v8 source class 5) ingestion path
  - DOCX (via python-docx)
  - HTML (via BeautifulSoup4)
  - Plain text (UTF-8, latin-1 fallback)

DR-9 note: PPTX + PDF table extraction use only permissively-licensed libraries
(python-pptx BSD, pdfplumber MIT). PyMuPDF / fitz is deliberately avoided
(AGPL-3 copyleft). Tables are flattened into the per-page/slide text stream so
the existing upload -> chunk -> NER -> entity-link pipeline (UserDocumentConnector)
picks them up with no further wiring, and are also surfaced structurally on
ExtractedDocument.tables.

Size limit configurable via MZ_DOC_UPLOAD_MAX_MB env var (default 25 MB).

The output ExtractedDocument is consumed by:
  - services/document_ner.py (LLM-based entity extraction)
  - connectors/user_document.py (UserDocumentConnector)
"""

from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass, field
from typing import Any

logger = logging.getLogger(__name__)


# ── Configuration ──────────────────────────────────────────────────

_DEFAULT_MAX_MB = 25
_PDF_MAGIC = b"%PDF-"


# ── Models ─────────────────────────────────────────────────────────

@dataclass
class ExtractedTable:
    """A table lifted from a document. ``rows`` is row-major cell text;
    ``page`` is the 0-based page (PDF) or slide (PPTX) index it came from."""
    rows: list[list[str]] = field(default_factory=list)
    page: int | None = None


@dataclass
class ExtractedDocument:
    """Result of document text extraction.

    Fields:
      format: detected file format ("pdf" | "pptx" | "docx" | "html" | "txt")
      pages: list of per-page/slide text strings (single element for non-paged)
      full_text: pages joined with "\\n\\n"
      tables: structured tables extracted from the document (DR-9)
      metadata: format-specific metadata (page_count, slide_count, encoding, ...)
    """
    format: str
    pages: list[str] = field(default_factory=list)
    full_text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    tables: list[ExtractedTable] = field(default_factory=list)


def _table_to_text(rows: list[list[str]]) -> str:
    """Flatten a table to a pipe-delimited text block so it joins the
    NER/embedding text stream. Drops fully-empty rows."""
    lines = []
    for row in rows:
        cells = [(c or "").strip() for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)


class UnsupportedFormatError(Exception):
    """Raised when the file format cannot be extracted."""


# ── Public API ─────────────────────────────────────────────────────

def extract_text(payload: bytes, filename: str = "") -> ExtractedDocument:
    """Extract text from an uploaded document.

    Args:
        payload: raw bytes of the document
        filename: original filename (used for format detection by extension)

    Returns:
        ExtractedDocument with per-page text + metadata.

    Raises:
        ValueError: if payload exceeds MZ_DOC_UPLOAD_MAX_MB
        UnsupportedFormatError: if format cannot be determined or extracted
    """
    # Size limit
    max_mb = int(os.getenv("MZ_DOC_UPLOAD_MAX_MB", str(_DEFAULT_MAX_MB)))
    size_mb = len(payload) / (1024 * 1024)
    if size_mb > max_mb:
        raise ValueError(
            f"document size {size_mb:.1f} MB exceeds limit of {max_mb} MB"
        )

    fmt = _detect_format(payload, filename)
    if fmt == "pdf":
        return _extract_pdf(payload)
    if fmt == "pptx":
        return _extract_pptx(payload)
    if fmt == "docx":
        return _extract_docx(payload)
    if fmt == "html":
        return _extract_html(payload)
    if fmt == "txt":
        return _extract_txt(payload)
    raise UnsupportedFormatError(
        f"cannot extract text from {filename!r}: unsupported or unknown format"
    )


# ── Format detection ───────────────────────────────────────────────

def _detect_format(payload: bytes, filename: str) -> str:
    """Detect format by extension first, falling back to magic bytes."""
    name = (filename or "").lower()
    # Extension-based (preferred — explicit user intent)
    if name.endswith(".pdf"):
        return "pdf"
    if name.endswith(".docx"):
        return "docx"
    if name.endswith(".pptx"):
        return "pptx"
    if name.endswith(".html") or name.endswith(".htm"):
        return "html"
    if name.endswith(".txt") or name.endswith(".md"):
        return "txt"

    # Magic-byte fallback for files with no/unknown extension
    if payload.startswith(_PDF_MAGIC):
        return "pdf"
    # DOCX is a ZIP — but ZIP magic alone isn't enough. Don't auto-detect DOCX
    # without an extension; treat it as unsupported to fail loudly rather than
    # mis-extract.

    # Try to decode as UTF-8 text — if it parses, treat as txt
    try:
        decoded = payload.decode("utf-8")
        # Heuristic: mostly printable
        if decoded and sum(1 for c in decoded if c.isprintable() or c in "\n\r\t") / max(len(decoded), 1) > 0.95:
            return "txt"
    except UnicodeDecodeError:
        pass

    return "unknown"


# ── Per-format extractors ──────────────────────────────────────────

def _extract_pdf(payload: bytes) -> ExtractedDocument:
    """Extract per-page text from a PDF using pdfplumber."""
    try:
        import pdfplumber
    except ImportError as exc:
        raise UnsupportedFormatError(
            "PDF extraction requires pdfplumber (pip install pdfplumber)"
        ) from exc

    pages: list[str] = []
    tables: list[ExtractedTable] = []
    try:
        with pdfplumber.open(io.BytesIO(payload)) as pdf:
            for idx, page in enumerate(pdf.pages):
                # extract_text returns None for empty pages — coerce to ""
                text = page.extract_text() or ""
                # DR-9: lift tables (10-K financials, trial-result tables) and
                # fold them into the text stream so NER/embedding sees them.
                page_tables = []
                try:
                    page_tables = page.extract_tables() or []
                except Exception:
                    logger.debug("table extraction failed on PDF page %d", idx)
                table_texts = []
                for raw in page_tables:
                    rows = [[(c or "").strip() for c in row] for row in raw]
                    tables.append(ExtractedTable(rows=rows, page=idx))
                    table_texts.append(_table_to_text(rows))
                if table_texts:
                    text = "\n\n".join([t for t in [text, *table_texts] if t])
                pages.append(text)
    except Exception as exc:
        # pdfplumber raises various low-level errors on corrupt PDFs.
        # Wrap in a clearer message so callers know the cause.
        raise ValueError(f"failed to parse PDF: {exc}") from exc

    full_text = "\n\n".join(pages)
    return ExtractedDocument(
        format="pdf",
        pages=pages,
        full_text=full_text,
        metadata={"page_count": len(pages), "table_count": len(tables)},
        tables=tables,
    )


def _extract_pptx(payload: bytes) -> ExtractedDocument:
    """Extract per-slide text, tables, and speaker notes from a PPTX deck.

    DR-9 / v8 source class 5 (conference decks). One ``page`` per slide. Tables
    are both flattened into the slide text and surfaced on ``tables``. Speaker
    notes are appended (prefixed) because conference decks often carry the
    quantitative detail in the notes."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise UnsupportedFormatError(
            "PPTX extraction requires python-pptx (pip install python-pptx)"
        ) from exc

    try:
        prs = Presentation(io.BytesIO(payload))
    except Exception as exc:
        raise ValueError(f"failed to parse PPTX: {exc}") from exc

    pages: list[str] = []
    tables: list[ExtractedTable] = []
    for idx, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows]
                tables.append(ExtractedTable(rows=rows, page=idx))
                parts.append(_table_to_text(rows))
            elif getattr(shape, "has_text_frame", False):
                txt = shape.text_frame.text.strip()
                if txt:
                    parts.append(txt)
        # Speaker notes — often where the data lives in a conference deck.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[notes] {notes}")
        except Exception:
            logger.debug("notes extraction failed on slide %d", idx)
        pages.append("\n".join(parts))

    full_text = "\n\n".join(pages)
    return ExtractedDocument(
        format="pptx",
        pages=pages,
        full_text=full_text,
        metadata={"slide_count": len(pages), "table_count": len(tables)},
        tables=tables,
    )


def _extract_docx(payload: bytes) -> ExtractedDocument:
    """Extract paragraph text from a DOCX file using python-docx."""
    try:
        from docx import Document
    except ImportError as exc:
        raise UnsupportedFormatError(
            "DOCX extraction requires python-docx"
        ) from exc

    try:
        doc = Document(io.BytesIO(payload))
    except Exception as exc:
        raise ValueError(f"failed to parse DOCX: {exc}") from exc

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)
    # DOCX has no native page concept — return single "page" with all text
    return ExtractedDocument(
        format="docx",
        pages=[full_text],
        full_text=full_text,
        metadata={"paragraph_count": len(paragraphs)},
    )


def _extract_html(payload: bytes) -> ExtractedDocument:
    """Extract visible text from HTML using BeautifulSoup4."""
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise UnsupportedFormatError(
            "HTML extraction requires beautifulsoup4"
        ) from exc

    # Encoding: BS4 will sniff
    soup = BeautifulSoup(payload, "html.parser")
    # Drop script/style content — never useful for NER
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ", strip=True)
    return ExtractedDocument(
        format="html",
        pages=[text],
        full_text=text,
        metadata={},
    )


def _extract_txt(payload: bytes) -> ExtractedDocument:
    """Plain-text decoding with utf-8 → latin-1 fallback."""
    if not payload:
        return ExtractedDocument(
            format="txt", pages=[], full_text="", metadata={"encoding": "utf-8"}
        )
    try:
        text = payload.decode("utf-8")
        encoding = "utf-8"
    except UnicodeDecodeError:
        text = payload.decode("latin-1")
        encoding = "latin-1"
    return ExtractedDocument(
        format="txt",
        pages=[text],
        full_text=text,
        metadata={"encoding": encoding},
    )
